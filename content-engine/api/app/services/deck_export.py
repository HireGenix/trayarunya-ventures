"""Branded deck exporters — turn a stored deck into a polished PPTX or PDF.

Both renderers honour the deck's resolved theme (brand primary/accent colours)
and lay each slide out per its layout so the export matches the on-screen design
intent: a designed, on-brand deck — not a blank template dump.
"""
from __future__ import annotations

import io
import logging
import urllib.request
from typing import Any

from app.models import Deck, DeckSlide

log = logging.getLogger("deck_export")

# 16:9 canvas in EMU-friendly inches for python-pptx.
_SLIDE_W = 13.333
_SLIDE_H = 7.5

_HERO = ("cover", "section", "cta", "image")
_SIDE = ("bullets", "agenda")


def _fetch_image(url: str | None, _cache: dict[str, bytes | None] = {}) -> bytes | None:
    """Download an image URL → bytes (best-effort, cached per process call)."""
    if not url or not isinstance(url, str) or not url.startswith("http"):
        return None
    if url in _cache:
        return _cache[url]
    data: bytes | None = None
    try:
        req = urllib.request.Request(url, headers={"User-Agent": "MarketIQ-Deck/1.0"})
        with urllib.request.urlopen(req, timeout=20) as resp:  # noqa: S310 — our own/Pexels URLs
            data = resp.read()
    except Exception as exc:  # noqa: BLE001 — imagery is best-effort
        log.warning("deck export image fetch failed (%s): %s", url, exc)
        data = None
    _cache[url] = data
    return data


def _hex_to_rgb(value: str | None, fallback: tuple[int, int, int]) -> tuple[int, int, int]:
    if not value or not isinstance(value, str):
        return fallback
    s = value.strip().lstrip("#")
    if len(s) == 3:
        s = "".join(c * 2 for c in s)
    if len(s) < 6:
        return fallback
    try:
        return int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16)
    except ValueError:
        return fallback


def _theme(deck: Deck) -> dict[str, Any]:
    t = deck.theme or {}
    from app.services.deck_themes import PPTX_FALLBACK, PDF_FONT_MAP
    heading_font = t.get("heading_font") or "Inter"
    body_font = t.get("body_font") or "Inter"
    return {
        "primary": _hex_to_rgb(t.get("primary"), (20, 187, 135)),
        "accent": _hex_to_rgb(t.get("accent"), (15, 168, 116)),
        "ink": _hex_to_rgb(t.get("ink"), (15, 23, 42)),
        "brand_name": t.get("brand_name") or "",
        "logo_url": t.get("logo_url") or "",
        "subtitle": (deck.meta or {}).get("subtitle", ""),
        "heading_font": heading_font,
        "body_font": body_font,
        "pptx_heading": PPTX_FALLBACK.get(heading_font, "Calibri"),
        "pptx_body": PPTX_FALLBACK.get(body_font, "Calibri"),
        "pdf_regular": PDF_FONT_MAP.get(body_font, ("Helvetica", "Helvetica-Bold"))[0],
        "pdf_bold": PDF_FONT_MAP.get(body_font, ("Helvetica", "Helvetica-Bold"))[1],
        "pdf_heading_bold": PDF_FONT_MAP.get(heading_font, ("Helvetica", "Helvetica-Bold"))[1],
    }


def _short_url(url: str, limit: int = 60) -> str:
    """Compact a URL to host + truncated path for footnotes/citations."""
    try:
        from urllib.parse import urlparse

        p = urlparse(url)
        host = p.netloc.lower()
        if host.startswith("www."):
            host = host[4:]
        s = host + (p.path if p.path and p.path != "/" else "")
    except Exception:  # noqa: BLE001
        s = url
    return s if len(s) <= limit else s[: limit - 1] + "…"


# Map controlled icon vocab → a unicode glyph that renders in PPTX/PDF fonts.
_ICON_GLYPH = {
    "brain": "✷", "shield": "❖", "chart": "▥", "bolt": "⚡", "check": "✓",
    "cost": "$", "users": "♚", "building": "▣", "rocket": "➤", "doc": "▤",
    "search": "⌕", "gear": "✦", "target": "◎", "idea": "✲", "globe": "◍",
    "lock": "▪", "clock": "◷", "star": "★", "graph": "▦", "layers": "▩",
    "flow": "➜", "money": "$",
}


def _glyph(icon: Any) -> str:
    key = str(icon or "").strip().lower()
    return _ICON_GLYPH.get(key, "•")


def _cell_glyph(value: str) -> str:
    """Render a matrix cell value as ✓ / ✗ / – or the short text itself."""
    v = str(value or "").strip().lower()
    if v in ("yes", "y", "true", "✓", "full", "1"):
        return "✓"
    if v in ("no", "n", "false", "✗", "none", "0", "-", "–"):
        return "✗"
    if v in ("partial", "some", "limited", "~"):
        return "~"
    return str(value or "")[:14]


def _callout_pptx(slide, d, add_rect, add_text, top, body_right, accent, ink,
                  RGBColor, PP_ALIGN, MSO_ANCHOR) -> None:
    text = d.get("callout")
    if not text or not isinstance(text, str):
        return
    tint = RGBColor(
        min(255, accent[0] + (255 - accent[0]) * 85 // 100),
        min(255, accent[1] + (255 - accent[1]) * 85 // 100),
        min(255, accent[2] + (255 - accent[2]) * 85 // 100),
    )
    h = 0.95
    y = _SLIDE_H - 0.75 - h
    add_rect(slide, 0.85, y, body_right - 0.85, h, tint)
    add_rect(slide, 0.85, y, 0.1, h, accent)
    add_text(slide, 1.15, y + 0.1, body_right - 1.4, h - 0.2,
             [("✲  " + text, 15, ink, True)], anchor=MSO_ANCHOR.MIDDLE)


def _callout_pdf(c, d, wrap, draw_lines, W, accent, ink, Color) -> None:
    text = d.get("callout")
    if not text or not isinstance(text, str):
        return
    tint = Color(
        accent.red + (1 - accent.red) * 0.85,
        accent.green + (1 - accent.green) * 0.85,
        accent.blue + (1 - accent.blue) * 0.85,
    )
    box_w = W - 180
    box_h = 88
    x, y = 90, 70
    c.setFillColor(tint)
    c.roundRect(x, y, box_w, box_h, 12, fill=1, stroke=0)
    c.setFillColor(accent)
    c.rect(x, y, 8, box_h, fill=1, stroke=0)
    lines = wrap("✲  " + text, "Helvetica-Bold", 16, box_w - 56)[:3]
    ty = y + box_h - 30
    draw_lines(x + 30, ty, lines, "Helvetica-Bold", 16, ink, leading=22)


# --------------------------------------------------------------------------- #
# PPTX
# --------------------------------------------------------------------------- #
def build_pptx(deck: Deck, slides: list[DeckSlide]) -> bytes:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
    from pptx.util import Inches, Pt

    th = _theme(deck)
    primary = RGBColor(*th["primary"])
    accent = RGBColor(*th["accent"])
    ink = RGBColor(*th["ink"])
    white = RGBColor(0xFF, 0xFF, 0xFF)
    muted = RGBColor(0x6B, 0x72, 0x80)

    prs = Presentation()
    prs.slide_width = Inches(_SLIDE_W)
    prs.slide_height = Inches(_SLIDE_H)
    blank = prs.slide_layouts[6]

    def add_rect(slide, x, y, w, h, color, line=False):
        from pptx.enum.shapes import MSO_SHAPE
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        if not line:
            shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def add_text(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, autofit=False):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        if autofit:
            # Let PowerPoint shrink the copy so long content never spills the box.
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        first = True
        for text, size, color, bold in runs:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = align
            p.space_after = Pt(6)
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = th.get("pptx_body", "Calibri")
        return box

    def footer(slide, dark=False):
        c = white if dark else muted
        add_text(slide, 0.5, _SLIDE_H - 0.5, 8, 0.35,
                 [(th["brand_name"], 10, c, False)], align=PP_ALIGN.LEFT)

    logo_bytes = _fetch_image(th.get("logo_url")) if th.get("logo_url") else None

    def add_logo(slide, dark=False):
        """Place the brand logo top-right (best-effort; sized to a 1.6in box)."""
        if not logo_bytes:
            return
        try:
            from PIL import Image  # noqa: PLC0415

            with Image.open(io.BytesIO(logo_bytes)) as im:
                iw, ih = im.size
            ratio = (iw / ih) if ih else 4.0
        except Exception:  # noqa: BLE001
            ratio = 4.0
        h = 0.6
        w = max(0.6, min(2.2, h * ratio))
        try:
            slide.shapes.add_picture(
                io.BytesIO(logo_bytes), Inches(_SLIDE_W - w - 0.6), Inches(0.45),
                Inches(w), Inches(h),
            )
        except Exception:  # noqa: BLE001
            pass

    def add_citation(slide, sources, dark=False):
        """Render a small 'Sources: …' footnote for cited content slides."""
        if not sources:
            return
        labels = []
        for src in sources[:3]:
            if isinstance(src, dict) and src.get("url"):
                labels.append(_short_url(src["url"], 46))
        if not labels:
            return
        c = white if dark else muted
        add_text(slide, _SLIDE_W - 7.3, _SLIDE_H - 0.5, 6.8, 0.35,
                 [("Sources: " + " · ".join(labels), 8, c, False)], align=PP_ALIGN.RIGHT)

    def add_image_fill(slide, x, y, w, h, data: bytes):
        """Place an image filling the box (stretched to box — fine for backdrops)."""
        try:
            slide.shapes.add_picture(io.BytesIO(data), Inches(x), Inches(y), Inches(w), Inches(h))
            return True
        except Exception:  # noqa: BLE001
            return False

    for s in slides:
        d = s.data or {}
        slide = prs.slides.add_slide(blank)
        layout = s.layout
        img = _fetch_image(d.get("image_url")) if layout in _HERO else None
        side_img = (
            _fetch_image(d.get("image_url")) if layout in _SIDE else None
        )

        if layout in _HERO:
            eyebrow = d.get("eyebrow") or ("" if layout != "cta" else "LET'S TALK")
            runs = []
            if eyebrow:
                runs.append((str(eyebrow).upper(), 16, accent, True))
            runs.append((d.get("title", ""), 44 if layout != "cover" else 48, white, True))
            sub = d.get("subtitle") or d.get("body") or ""
            if sub:
                runs.append((sub, 22, white, False))
            if layout == "cta" and d.get("cta"):
                runs.append(("➜  " + d["cta"], 20, accent, True))

            if img and layout == "image":
                # Full-bleed image with a solid brand band at the bottom for the text.
                add_image_fill(slide, 0, 0, _SLIDE_W, _SLIDE_H, img)
                band_h = 2.5
                add_rect(slide, 0, _SLIDE_H - band_h, _SLIDE_W, band_h, primary)
                add_rect(slide, 0, _SLIDE_H - band_h, _SLIDE_W, 0.12, accent)
                add_text(slide, 0.9, _SLIDE_H - band_h + 0.3, _SLIDE_W - 1.8, band_h - 0.5,
                         runs, anchor=MSO_ANCHOR.MIDDLE)
            elif img:
                # Split: brand panel on the left, image on the right.
                panel_w = _SLIDE_W * 0.6
                add_rect(slide, 0, 0, panel_w, _SLIDE_H, primary)
                add_rect(slide, 0, _SLIDE_H - 0.35, panel_w, 0.35, accent)
                add_image_fill(slide, panel_w, 0, _SLIDE_W - panel_w, _SLIDE_H, img)
                add_text(slide, 0.9, 1.4, panel_w - 1.5, _SLIDE_H - 2.4, runs,
                         anchor=MSO_ANCHOR.MIDDLE)
            else:
                # No image — full-bleed brand panel (original behaviour).
                add_rect(slide, 0, 0, _SLIDE_W, _SLIDE_H, primary)
                add_rect(slide, 0, _SLIDE_H - 0.35, _SLIDE_W, 0.35, accent)
                add_text(slide, 0.9, 1.4, _SLIDE_W - 1.8, _SLIDE_H - 2.4, runs,
                         anchor=MSO_ANCHOR.MIDDLE)
            footer(slide, dark=True)
            add_logo(slide, dark=True)
            add_citation(slide, d.get("sources"), dark=True)
            _notes(slide, s)
            continue

        # Light content slides — accent rail + title.
        add_rect(slide, 0, 0, 0.28, _SLIDE_H, primary)
        # Optional tall side photo on bullets/agenda → content keeps the left area.
        body_right = _SLIDE_W - 1.4
        if side_img is not None:
            panel_w = _SLIDE_W * 0.32
            panel_x = _SLIDE_W - panel_w - 0.45
            add_image_fill(slide, panel_x, 0.7, panel_w, _SLIDE_H - 1.4, side_img)
            add_rect(slide, panel_x, _SLIDE_H - 0.7, panel_w, 0.12, accent)
            body_right = panel_x - 0.4
        title = d.get("title", "")
        if title:
            add_text(slide, 0.8, 0.6, body_right - 0.8, 1.1,
                     [(title, 32, ink, True)])
        if d.get("subtitle"):
            add_text(slide, 0.8, 1.55, body_right - 0.8, 0.6,
                     [(d["subtitle"], 18, muted, False)])

        top = 2.3
        if layout == "agenda":
            for i, item in enumerate(d.get("items", [])):
                y = top + i * 0.62
                add_rect(slide, 0.85, y + 0.06, 0.3, 0.3, accent)
                add_text(slide, 1.3, y, body_right - 1.3, 0.55, [(str(item), 20, ink, False)], autofit=True)
        elif layout == "bullets":
            for i, b in enumerate(d.get("bullets", [])):
                y = top + i * 0.78
                add_rect(slide, 0.85, y + 0.12, 0.18, 0.18, primary)
                runs = [(b.get("heading", ""), 20, ink, True)]
                if b.get("body"):
                    runs.append((b["body"], 16, muted, False))
                add_text(slide, 1.25, y, body_right - 1.25, 0.8, runs, autofit=True)
        elif layout in ("two_column", "comparison"):
            for ci, key in enumerate(("left", "right")):
                col = d.get(key) or {}
                x = 0.85 + ci * ((_SLIDE_W - 1.7) / 2 + 0.2)
                w = (_SLIDE_W - 1.7) / 2 - 0.2
                add_rect(slide, x, top, w, _SLIDE_H - top - 0.7,
                         RGBColor(0xF3, 0xF4, 0xF6))
                add_rect(slide, x, top, w, 0.12, accent if ci else primary)
                runs = [(col.get("heading", ""), 20, ink, True)]
                if col.get("body"):
                    runs.append((col["body"], 16, muted, False))
                for it in col.get("items", []) or []:
                    runs.append(("•  " + str(it), 16, ink, False))
                add_text(slide, x + 0.3, top + 0.3, w - 0.6, _SLIDE_H - top - 1.2, runs, autofit=True)
        elif layout == "stats":
            stats = d.get("stats", [])[:4]
            n = max(1, len(stats))
            gap = 0.3
            w = (_SLIDE_W - 1.7 - gap * (n - 1)) / n
            for i, st in enumerate(stats):
                x = 0.85 + i * (w + gap)
                add_rect(slide, x, top, w, 2.6, RGBColor(0xF3, 0xF4, 0xF6))
                add_rect(slide, x, top, w, 0.14, primary)
                add_text(slide, x, top + 0.5, w, 1.2, [(st.get("value", ""), 46, primary, True)],
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
                add_text(slide, x + 0.15, top + 1.7, w - 0.3, 0.8, [(st.get("label", ""), 15, muted, False)],
                         align=PP_ALIGN.CENTER)
        elif layout == "quote":
            add_text(slide, 1.2, 2.0, _SLIDE_W - 2.4, 3.0,
                     [("“" + d.get("quote", "") + "”", 30, ink, True)], anchor=MSO_ANCHOR.MIDDLE)
            if d.get("attribution"):
                add_text(slide, 1.2, 5.0, _SLIDE_W - 2.4, 0.6,
                         [("— " + d["attribution"], 18, accent, True)])
        elif layout == "timeline":
            steps = d.get("steps", [])[:6]
            n = max(1, len(steps))
            w = (_SLIDE_W - 1.7) / n
            for i, stp in enumerate(steps):
                x = 0.85 + i * w
                add_rect(slide, x, top + 0.4, w - 0.25, 0.1, accent)
                add_rect(slide, x, top + 0.18, 0.5, 0.5, primary)
                add_text(slide, x, top + 0.22, 0.5, 0.45, [(str(i + 1), 18, white, True)],
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
                runs = [(stp.get("label", ""), 17, ink, True)]
                if stp.get("body"):
                    runs.append((stp["body"], 14, muted, False))
                add_text(slide, x, top + 0.9, w - 0.2, 2.4, runs, autofit=True)
        elif layout == "references":
            items = d.get("items", []) or []
            half = (len(items) + 1) // 2
            cols = [items[:half], items[half:]] if len(items) > 6 else [items, []]
            for ci, col in enumerate(cols):
                x = 0.85 + ci * ((_SLIDE_W - 1.7) / 2)
                runs = []
                for it in col:
                    if not isinstance(it, dict):
                        continue
                    label = str(it.get("label") or it.get("url") or "")[:70]
                    url = _short_url(str(it.get("url") or ""), 60)
                    runs.append(("•  " + label, 13, ink, True))
                    if url:
                        runs.append(("    " + url, 11, accent, False))
                if runs:
                    add_text(slide, x, top, (_SLIDE_W - 1.7) / 2 - 0.2, _SLIDE_H - top - 0.7, runs)
        elif layout == "cards":
            cards = (d.get("cards") or [])[:6]
            n = max(1, len(cards))
            cols = 2 if n <= 4 else 3
            rows = (n + cols - 1) // cols
            gap = 0.3
            avail_h = _SLIDE_H - top - 0.9
            cw = (body_right - 0.85 - gap * (cols - 1)) / cols
            ch = (avail_h - gap * (rows - 1)) / rows
            for i, card in enumerate(cards):
                r, cc = divmod(i, cols)
                x = 0.85 + cc * (cw + gap)
                y = top + r * (ch + gap)
                add_rect(slide, x, y, cw, ch, RGBColor(0xF3, 0xF4, 0xF6))
                add_rect(slide, x, y, cw, 0.12, primary)
                runs = [(_glyph(card.get("icon")) + "  " + str(card.get("heading", "")), 18, ink, True)]
                if card.get("body"):
                    runs.append((str(card["body"]), 14, muted, False))
                add_text(slide, x + 0.28, y + 0.28, cw - 0.56, ch - 0.5, runs, autofit=True)
        elif layout == "process":
            steps = (d.get("steps") or [])[:6]
            n = max(1, len(steps))
            cols = n if n <= 3 else (n + 1) // 2
            rows = 1 if n <= 3 else 2
            gap = 0.3
            avail_h = _SLIDE_H - top - 0.9
            cw = (body_right - 0.85 - gap * (cols - 1)) / cols
            ch = (avail_h - gap * (rows - 1)) / rows
            for i, stp in enumerate(steps):
                r, cc = divmod(i, cols)
                x = 0.85 + cc * (cw + gap)
                y = top + r * (ch + gap)
                add_rect(slide, x, y, cw, ch, RGBColor(0xF3, 0xF4, 0xF6))
                add_rect(slide, x, y, 0.5, 0.5, primary)
                add_text(slide, x, y + 0.04, 0.5, 0.45, [(str(i + 1), 18, white, True)],
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
                runs = [(_glyph(stp.get("icon")) + "  " + str(stp.get("heading", "")), 17, ink, True)]
                if stp.get("body"):
                    runs.append((str(stp["body"]), 14, muted, False))
                add_text(slide, x + 0.28, y + 0.62, cw - 0.56, ch - 0.8, runs, autofit=True)
        elif layout == "comparison_matrix":
            columns = [str(c) for c in (d.get("columns") or [])]
            mrows = (d.get("rows") or [])[:8]
            ncol = max(1, len(columns))
            label_w = (body_right - 0.85) * 0.34
            cell_w = (body_right - 0.85 - label_w) / max(1, ncol - 1)
            row_h = min(0.62, (_SLIDE_H - top - 0.9) / max(1, len(mrows) + 1))
            # header
            for ci, ctext in enumerate(columns):
                x = 0.85 + (0 if ci == 0 else label_w + (ci - 1) * cell_w)
                w = label_w if ci == 0 else cell_w
                add_text(slide, x + 0.1, top, w - 0.2, row_h,
                         [(ctext, 14, primary if ci == 1 else ink, True)],
                         align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER,
                         anchor=MSO_ANCHOR.MIDDLE)
            add_rect(slide, 0.85, top + row_h, body_right - 0.85, 0.03, primary)
            for ri, row in enumerate(mrows):
                y = top + row_h * (ri + 1) + 0.06
                if ri == 0:
                    add_rect(slide, 0.85, y, body_right - 0.85, row_h, RGBColor(0xEC, 0xFD, 0xF5))
                add_text(slide, 0.95, y, label_w - 0.2, row_h,
                         [(str(row.get("label", "")), 14, ink, ri == 0)],
                         anchor=MSO_ANCHOR.MIDDLE)
                cells = [str(x) for x in (row.get("cells") or [])]
                for ci in range(ncol - 1):
                    x = 0.85 + label_w + ci * cell_w
                    val = _cell_glyph(cells[ci] if ci < len(cells) else "")
                    add_text(slide, x, y, cell_w, row_h, [(val, 16, ink, True)],
                             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        elif layout == "chart":
            _add_pptx_chart(slide, d, top, body_right, th, prs, RGBColor, Inches, Pt, PP_ALIGN)

        _callout_pptx(slide, d, add_rect, add_text, top, body_right, accent, ink,
                      RGBColor, PP_ALIGN, MSO_ANCHOR)
        footer(slide)
        add_logo(slide)
        if layout != "references":
            add_citation(slide, d.get("sources"))
        _notes(slide, s)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _notes(slide, s: DeckSlide) -> None:
    if s.speaker_notes:
        try:
            slide.notes_slide.notes_text_frame.text = s.speaker_notes
        except Exception:  # noqa: BLE001
            pass


def _add_pptx_chart(slide, d, top, body_right, th, prs, RGBColor, Inches, Pt, PP_ALIGN):
    """Render a native PPTX chart (bar/line/pie/area/column) for a 'chart' slide."""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    chart_type = (d.get("chart_type") or "bar").lower()
    series_data = d.get("series") or []
    labels = d.get("labels") or []

    if not series_data:
        return

    type_map = {
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "area": XL_CHART_TYPE.AREA,
        "pie": XL_CHART_TYPE.PIE,
    }
    xl_type = type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

    chart_data = CategoryChartData()
    chart_data.categories = labels if labels else [f"C{i+1}" for i in range(max(len(s.get("values", [])) for s in series_data))]

    for sr in series_data[:6]:
        name = sr.get("name") or "Series"
        values = sr.get("values") or []
        chart_data.add_series(name, values)

    x = Inches(0.85)
    y = Inches(top + 0.2)
    w = Inches(body_right - 0.85)
    h = Inches(_SLIDE_H - top - 1.4)

    chart_frame = slide.shapes.add_chart(xl_type, x, y, w, h, chart_data)
    chart = chart_frame.chart
    chart.has_legend = len(series_data) > 1

    # Style the chart with theme colours
    primary_rgb = th["primary"]
    accent_rgb = th["accent"]
    colours = [
        RGBColor(*primary_rgb),
        RGBColor(*accent_rgb),
        RGBColor(107, 114, 128),
        RGBColor(59, 130, 246),
        RGBColor(245, 158, 11),
        RGBColor(239, 68, 68),
    ]
    plot = chart.plots[0]
    for i, sr in enumerate(plot.series):
        fill = sr.format.fill
        fill.solid()
        fill.fore_color.rgb = colours[i % len(colours)]


def _draw_pdf_chart(c, d, content_top, text_w, primary, accent, ink, muted, Color, wrap, draw_lines):
    """Draw a simple chart in the PDF export using reportlab drawing primitives."""
    chart_type = (d.get("chart_type") or "bar").lower()
    series_data = d.get("series") or []
    labels = d.get("labels") or []

    if not series_data:
        return

    all_values = []
    for sr in series_data:
        all_values.extend(sr.get("values") or [])
    if not all_values:
        return

    max_val = max(abs(v) for v in all_values) if all_values else 1
    if max_val == 0:
        max_val = 1

    chart_x = 90
    chart_w = text_w - 40
    chart_h = min(320, content_top - 120)
    chart_y = content_top - chart_h - 30

    colours = [primary, accent, Color(0.42, 0.45, 0.50), Color(0.23, 0.51, 0.96),
               Color(0.96, 0.62, 0.04), Color(0.94, 0.27, 0.27)]

    if chart_type == "pie" and series_data:
        # Simple pie chart
        vals = series_data[0].get("values") or []
        total = sum(abs(v) for v in vals) or 1
        cx, cy = chart_x + chart_w / 2, chart_y + chart_h / 2
        radius = min(chart_w, chart_h) / 2.5
        start_angle = 90
        for i, v in enumerate(vals[:8]):
            sweep = (abs(v) / total) * 360
            c.setFillColor(colours[i % len(colours)])
            c.wedge(cx - radius, cy - radius, cx + radius, cy + radius,
                    start_angle, sweep, fill=1, stroke=1)
            start_angle += sweep
        # Labels
        label_y = chart_y - 10
        for i, lbl in enumerate(labels[:len(vals)]):
            c.setFillColor(colours[i % len(colours)])
            c.rect(chart_x + i * 120, label_y, 10, 10, fill=1, stroke=0)
            c.setFillColor(ink)
            c.setFont("Helvetica", 11)
            c.drawString(chart_x + i * 120 + 15, label_y + 1, str(lbl)[:20])
    else:
        # Bar / column / line / area
        n_cats = max(len(labels), max((len(s.get("values", [])) for s in series_data), default=1))
        bar_group_w = chart_w / max(1, n_cats)
        n_series = len(series_data)
        bar_w = bar_group_w / (n_series + 1)

        # Axis
        c.setStrokeColor(Color(0.8, 0.8, 0.8))
        c.line(chart_x, chart_y, chart_x, chart_y + chart_h)
        c.line(chart_x, chart_y, chart_x + chart_w, chart_y)

        for si, sr in enumerate(series_data[:6]):
            vals = sr.get("values") or []
            colour = colours[si % len(colours)]
            c.setFillColor(colour)
            c.setStrokeColor(colour)

            if chart_type in ("line", "area"):
                points = []
                for i, v in enumerate(vals[:n_cats]):
                    px = chart_x + (i + 0.5) * bar_group_w
                    py = chart_y + (v / max_val) * chart_h
                    points.append((px, py))
                if chart_type == "area" and points:
                    c.setFillColor(Color(colour.red, colour.green, colour.blue, alpha=0.3))
                    path = c.beginPath()
                    path.moveTo(points[0][0], chart_y)
                    for px, py in points:
                        path.lineTo(px, py)
                    path.lineTo(points[-1][0], chart_y)
                    path.close()
                    c.drawPath(path, fill=1, stroke=0)
                    c.setStrokeColor(colour)
                for j in range(len(points) - 1):
                    c.line(points[j][0], points[j][1], points[j+1][0], points[j+1][1])
                for px, py in points:
                    c.circle(px, py, 3, fill=1, stroke=0)
            else:
                for i, v in enumerate(vals[:n_cats]):
                    bx = chart_x + i * bar_group_w + si * bar_w + bar_w * 0.2
                    bh = (v / max_val) * chart_h
                    c.rect(bx, chart_y, bar_w * 0.8, max(1, bh), fill=1, stroke=0)

        # Category labels
        c.setFillColor(ink)
        c.setFont("Helvetica", 10)
        for i, lbl in enumerate(labels[:n_cats]):
            lx = chart_x + (i + 0.5) * bar_group_w
            c.drawCentredString(lx, chart_y - 18, str(lbl)[:15])

        # Legend
        if n_series > 1:
            ly = chart_y - 40
            for si, sr in enumerate(series_data[:6]):
                lx = chart_x + si * 130
                c.setFillColor(colours[si % len(colours)])
                c.rect(lx, ly, 10, 10, fill=1, stroke=0)
                c.setFillColor(ink)
                c.setFont("Helvetica", 10)
                c.drawString(lx + 15, ly + 1, str(sr.get("name", ""))[:18])


# --------------------------------------------------------------------------- #
# PDF
# --------------------------------------------------------------------------- #
def build_pdf(deck: Deck, slides: list[DeckSlide]) -> bytes:
    from reportlab.lib.colors import Color
    from reportlab.lib.utils import ImageReader, simpleSplit
    from reportlab.pdfgen import canvas

    th = _theme(deck)
    W, H = 1280, 720  # 16:9 points

    def col(rgb, a=1.0):
        return Color(rgb[0] / 255, rgb[1] / 255, rgb[2] / 255, alpha=a)

    primary, accent, ink = col(th["primary"]), col(th["accent"]), col(th["ink"])
    muted = Color(0.42, 0.45, 0.50)
    panel = Color(0.95, 0.96, 0.97)
    white = Color(1, 1, 1)
    hd_bold = th.get("pdf_heading_bold", "Helvetica-Bold")
    bd_reg = th.get("pdf_regular", "Helvetica")
    bd_bold = th.get("pdf_bold", "Helvetica-Bold")

    buf = io.BytesIO()
    c = canvas.Canvas(buf, pagesize=(W, H))

    def wrap(text, font, size, max_w):
        return simpleSplit(str(text or ""), font, size, max_w)

    def img_reader(data: bytes | None):
        if not data:
            return None
        try:
            return ImageReader(io.BytesIO(data))
        except Exception:  # noqa: BLE001
            return None

    def draw_lines(x, y, lines, font, size, color, leading=None, align="left", max_w=W):
        leading = leading or size * 1.25
        c.setFont(font, size)
        c.setFillColor(color)
        for ln in lines:
            if align == "center":
                c.drawCentredString(x, y, ln)
            else:
                c.drawString(x, y, ln)
            y -= leading
        return y

    logo_reader = img_reader(_fetch_image(th.get("logo_url"))) if th.get("logo_url") else None

    def draw_logo():
        if logo_reader is None:
            return
        try:
            iw, ih = logo_reader.getSize()
            ratio = (iw / ih) if ih else 4.0
        except Exception:  # noqa: BLE001
            ratio = 4.0
        h = 54
        w = max(54, min(200, h * ratio))
        try:
            c.drawImage(logo_reader, W - w - 56, H - h - 40, w, h,
                        preserveAspectRatio=True, mask="auto")
        except Exception:  # noqa: BLE001
            pass

    def draw_citation(sources, dark=False):
        if not sources:
            return
        labels = [
            _short_url(s["url"], 42)
            for s in sources[:3]
            if isinstance(s, dict) and s.get("url")
        ]
        if not labels:
            return
        c.setFont("Helvetica", 11)
        c.setFillColor(white if dark else muted)
        c.drawRightString(W - 80, 36, "Sources: " + " · ".join(labels))

    for s in slides:
        d = s.data or {}
        layout = s.layout
        hero_img = img_reader(_fetch_image(d.get("image_url"))) if layout in _HERO else None

        if layout in _HERO:
            eyebrow = d.get("eyebrow") or ("LET'S TALK" if layout == "cta" else "")
            sub = d.get("subtitle") or d.get("body") or ""

            if hero_img is not None and layout == "image":
                # Full-bleed image + solid brand band at the bottom for the text.
                c.drawImage(hero_img, 0, 0, W, H, preserveAspectRatio=False, mask="auto")
                band_h = 230
                c.setFillColor(primary)
                c.rect(0, 0, W, band_h, fill=1, stroke=0)
                c.setFillColor(accent)
                c.rect(0, band_h - 10, W, 10, fill=1, stroke=0)
                ty = band_h - 70
                if eyebrow:
                    c.setFont("Helvetica-Bold", 20); c.setFillColor(accent)
                    c.drawString(80, band_h - 36, str(eyebrow).upper())
                for ln in wrap(d.get("title", ""), "Helvetica-Bold", 44, W - 160):
                    c.setFont("Helvetica-Bold", 44); c.setFillColor(white)
                    c.drawString(80, ty, ln); ty -= 52
                if sub:
                    for ln in wrap(sub, "Helvetica", 22, W - 160):
                        c.setFont("Helvetica", 22); c.setFillColor(white)
                        c.drawString(80, ty, ln); ty -= 30
            else:
                panel_w = W * 0.6 if hero_img is not None else W
                if hero_img is not None:
                    c.drawImage(hero_img, panel_w, 0, W - panel_w, H,
                                preserveAspectRatio=False, mask="auto")
                c.setFillColor(primary)
                c.rect(0, 0, panel_w, H, fill=1, stroke=0)
                c.setFillColor(accent)
                c.rect(0, 0, panel_w, 26, fill=1, stroke=0)
                tx, tw = 80, panel_w - 160
                y = H - 230
                if eyebrow:
                    c.setFont("Helvetica-Bold", 22); c.setFillColor(accent)
                    c.drawString(tx, y, str(eyebrow).upper()); y -= 50
                for ln in wrap(d.get("title", ""), "Helvetica-Bold", 56 if hero_img is None else 46, tw):
                    c.setFont("Helvetica-Bold", 56 if hero_img is None else 46)
                    c.setFillColor(white); c.drawString(tx, y, ln); y -= 64
                if sub:
                    y -= 10
                    for ln in wrap(sub, "Helvetica", 26, tw):
                        c.setFont("Helvetica", 26); c.setFillColor(white)
                        c.drawString(tx, y, ln); y -= 36
                if layout == "cta" and d.get("cta"):
                    c.setFont("Helvetica-Bold", 26); c.setFillColor(accent)
                    c.drawString(tx, y - 10, "\u2794  " + d["cta"])
                c.setFont("Helvetica", 16); c.setFillColor(white)
                c.drawString(tx, 40, th["brand_name"])
            draw_logo()
            draw_citation(d.get("sources"), dark=True)
            c.showPage()
            continue

        # Light content slide.
        c.setFillColor(white)
        c.rect(0, 0, W, H, fill=1, stroke=0)
        c.setFillColor(primary)
        c.rect(0, 0, 26, H, fill=1, stroke=0)

        # Optional tall side photo on bullets/agenda.
        side_reader = (
            img_reader(_fetch_image(d.get("image_url"))) if layout in _SIDE else None
        )
        text_w = W - 220
        if side_reader is not None:
            panel_w = W * 0.32
            panel_x = W - panel_w - 36
            c.drawImage(side_reader, panel_x, 60, panel_w, H - 120,
                        preserveAspectRatio=False, mask="auto")
            c.setFillColor(accent)
            c.rect(panel_x, 60, panel_w, 8, fill=1, stroke=0)
            text_w = panel_x - 90 - 24

        y = H - 110
        title_lines = wrap(d.get("title", ""), "Helvetica-Bold", 40, text_w)
        y = draw_lines(90, y, title_lines, "Helvetica-Bold", 40, ink, leading=48)
        if d.get("subtitle"):
            y = draw_lines(90, y - 6, wrap(d["subtitle"], "Helvetica", 22, text_w),
                           "Helvetica", 22, muted, leading=28)
        y -= 24
        content_top = y

        if layout == "agenda":
            for item in d.get("items", []):
                c.setFillColor(accent)
                c.rect(92, y - 4, 16, 16, fill=1, stroke=0)
                draw_lines(124, y, wrap(item, "Helvetica", 24, text_w - 40),
                           "Helvetica", 24, ink, leading=30)
                y -= 52
        elif layout == "bullets":
            for b in d.get("bullets", []):
                c.setFillColor(primary)
                c.circle(100, y + 4, 7, fill=1, stroke=0)
                yy = draw_lines(124, y, wrap(b.get("heading", ""), "Helvetica-Bold", 24, text_w - 40),
                                "Helvetica-Bold", 24, ink, leading=30)
                if b.get("body"):
                    yy = draw_lines(124, yy, wrap(b["body"], "Helvetica", 18, text_w - 40),
                                    "Helvetica", 18, muted, leading=24)
                y = yy - 18
        elif layout in ("two_column", "comparison"):
            cw = (W - 220) / 2
            for ci, key in enumerate(("left", "right")):
                cold = d.get(key) or {}
                x = 90 + ci * (cw + 40)
                c.setFillColor(panel)
                c.roundRect(x, 120, cw, content_top - 120, 14, fill=1, stroke=0)
                c.setFillColor(accent if ci else primary)
                c.rect(x, content_top - 14, cw, 10, fill=1, stroke=0)
                yy = content_top - 50
                yy = draw_lines(x + 28, yy, wrap(cold.get("heading", ""), "Helvetica-Bold", 24, cw - 56),
                                "Helvetica-Bold", 24, ink, leading=30)
                if cold.get("body"):
                    yy = draw_lines(x + 28, yy - 4, wrap(cold["body"], "Helvetica", 18, cw - 56),
                                    "Helvetica", 18, muted, leading=24)
                for it in cold.get("items", []) or []:
                    yy = draw_lines(x + 28, yy - 2, wrap("•  " + str(it), "Helvetica", 18, cw - 56),
                                    "Helvetica", 18, ink, leading=24)
        elif layout == "stats":
            stats = d.get("stats", [])[:4]
            n = max(1, len(stats))
            gap = 30
            cw = (W - 180 - gap * (n - 1)) / n
            for i, st in enumerate(stats):
                x = 90 + i * (cw + gap)
                c.setFillColor(panel)
                c.roundRect(x, content_top - 250, cw, 240, 14, fill=1, stroke=0)
                c.setFillColor(primary)
                c.rect(x, content_top - 22, cw, 12, fill=1, stroke=0)
                c.setFillColor(primary)
                c.setFont("Helvetica-Bold", 64)
                c.drawCentredString(x + cw / 2, content_top - 120, str(st.get("value", "")))
                c.setFillColor(muted)
                for j, ln in enumerate(wrap(st.get("label", ""), "Helvetica", 18, cw - 40)):
                    c.setFont("Helvetica", 18)
                    c.drawCentredString(x + cw / 2, content_top - 160 - j * 24, ln)
        elif layout == "quote":
            c.setFillColor(accent)
            c.setFont("Helvetica-Bold", 120)
            c.drawString(80, H / 2 + 40, "\u201C")
            yy = H / 2 + 60
            yy = draw_lines(W / 2, yy, wrap(d.get("quote", ""), "Helvetica-Bold", 36, W - 320),
                            "Helvetica-Bold", 36, ink, leading=46, align="center", max_w=W - 320)
            if d.get("attribution"):
                c.setFillColor(accent)
                c.setFont("Helvetica-Bold", 22)
                c.drawCentredString(W / 2, yy - 10, "\u2014 " + d["attribution"])
        elif layout == "timeline":
            steps = d.get("steps", [])[:6]
            n = max(1, len(steps))
            cw = (W - 180) / n
            for i, stp in enumerate(steps):
                x = 90 + i * cw
                c.setFillColor(accent)
                c.rect(x + 30, content_top - 40, cw - 30, 6, fill=1, stroke=0)
                c.setFillColor(primary)
                c.circle(x + 24, content_top - 36, 22, fill=1, stroke=0)
                c.setFillColor(white)
                c.setFont("Helvetica-Bold", 22)
                c.drawCentredString(x + 24, content_top - 44, str(i + 1))
                yy = content_top - 90
                yy = draw_lines(x, yy, wrap(stp.get("label", ""), "Helvetica-Bold", 20, cw - 30),
                                "Helvetica-Bold", 20, ink, leading=26)
                if stp.get("body"):
                    draw_lines(x, yy, wrap(stp["body"], "Helvetica", 16, cw - 30),
                               "Helvetica", 16, muted, leading=21)
        elif layout == "references":
            items = d.get("items", []) or []
            half = (len(items) + 1) // 2
            cols = [items[:half], items[half:]] if len(items) > 6 else [items, []]
            cw = (W - 220) / 2
            for ci, colitems in enumerate(cols):
                x = 90 + ci * (cw + 40)
                yy = content_top
                for it in colitems:
                    if not isinstance(it, dict):
                        continue
                    label = str(it.get("label") or it.get("url") or "")
                    url = _short_url(str(it.get("url") or ""), 56)
                    yy = draw_lines(x, yy, wrap("•  " + label, "Helvetica-Bold", 16, cw - 20),
                                    "Helvetica-Bold", 16, ink, leading=21)
                    if url:
                        yy = draw_lines(x + 18, yy, [url], "Helvetica", 13, accent, leading=20)
                    yy -= 8
        elif layout == "cards":
            cards = (d.get("cards") or [])[:6]
            n = max(1, len(cards))
            cols = 2 if n <= 4 else 3
            rows = (n + cols - 1) // cols
            gap = 28
            avail_h = content_top - 90
            cw = (W - 180 - gap * (cols - 1)) / cols
            ch = (avail_h - gap * (rows - 1)) / rows
            for i, card in enumerate(cards):
                r, cc = divmod(i, cols)
                x = 90 + cc * (cw + gap)
                ytop = content_top - r * (ch + gap)
                c.setFillColor(panel)
                c.roundRect(x, ytop - ch, cw, ch, 14, fill=1, stroke=0)
                c.setFillColor(primary)
                c.rect(x, ytop - 10, cw, 10, fill=1, stroke=0)
                yy = ytop - 40
                yy = draw_lines(x + 24, yy, wrap(_glyph(card.get("icon")) + "  " + str(card.get("heading", "")),
                                "Helvetica-Bold", 20, cw - 48), "Helvetica-Bold", 20, ink, leading=26)
                if card.get("body"):
                    draw_lines(x + 24, yy - 2, wrap(str(card["body"]), "Helvetica", 15, cw - 48),
                               "Helvetica", 15, muted, leading=20)
        elif layout == "process":
            steps = (d.get("steps") or [])[:6]
            n = max(1, len(steps))
            cols = n if n <= 3 else (n + 1) // 2
            rows = 1 if n <= 3 else 2
            gap = 28
            avail_h = content_top - 90
            cw = (W - 180 - gap * (cols - 1)) / cols
            ch = (avail_h - gap * (rows - 1)) / rows
            for i, stp in enumerate(steps):
                r, cc = divmod(i, cols)
                x = 90 + cc * (cw + gap)
                ytop = content_top - r * (ch + gap)
                c.setFillColor(panel)
                c.roundRect(x, ytop - ch, cw, ch, 14, fill=1, stroke=0)
                c.setFillColor(primary)
                c.circle(x + 34, ytop - 34, 20, fill=1, stroke=0)
                c.setFillColor(white)
                c.setFont("Helvetica-Bold", 20)
                c.drawCentredString(x + 34, ytop - 41, str(i + 1))
                yy = ytop - 80
                yy = draw_lines(x + 24, yy, wrap(_glyph(stp.get("icon")) + "  " + str(stp.get("heading", "")),
                                "Helvetica-Bold", 19, cw - 48), "Helvetica-Bold", 19, ink, leading=24)
                if stp.get("body"):
                    draw_lines(x + 24, yy - 2, wrap(str(stp["body"]), "Helvetica", 15, cw - 48),
                               "Helvetica", 15, muted, leading=20)
        elif layout == "comparison_matrix":
            columns = [str(x) for x in (d.get("columns") or [])]
            mrows = (d.get("rows") or [])[:8]
            ncol = max(1, len(columns))
            label_w = (W - 180) * 0.34
            cell_w = (W - 180 - label_w) / max(1, ncol - 1)
            row_h = min(54, (content_top - 70) / max(1, len(mrows) + 1))
            hy = content_top
            for ci, ctext in enumerate(columns):
                if ci == 0:
                    c.setFont("Helvetica-Bold", 16); c.setFillColor(ink)
                    c.drawString(94, hy - 18, ctext[:30])
                else:
                    cx = 90 + label_w + (ci - 1) * cell_w + cell_w / 2
                    c.setFont("Helvetica-Bold", 15)
                    c.setFillColor(primary if ci == 1 else ink)
                    c.drawCentredString(cx, hy - 18, ctext[:18])
            c.setFillColor(primary)
            c.rect(90, hy - row_h + 4, W - 180, 2, fill=1, stroke=0)
            for ri, row in enumerate(mrows):
                ry = hy - row_h * (ri + 1)
                if ri == 0:
                    c.setFillColor(Color(0.92, 0.99, 0.96))
                    c.rect(90, ry - 4, W - 180, row_h, fill=1, stroke=0)
                c.setFont("Helvetica-Bold" if ri == 0 else "Helvetica", 15)
                c.setFillColor(ink)
                c.drawString(98, ry + row_h / 2 - 8, str(row.get("label", ""))[:34])
                cells = [str(x) for x in (row.get("cells") or [])]
                for ci in range(ncol - 1):
                    cx = 90 + label_w + ci * cell_w + cell_w / 2
                    val = _cell_glyph(cells[ci] if ci < len(cells) else "")
                    c.setFont("Helvetica-Bold", 17); c.setFillColor(ink)
                    c.drawCentredString(cx, ry + row_h / 2 - 8, val)

        elif layout == "chart":
            _draw_pdf_chart(c, d, content_top, text_w, primary, accent, ink, muted, Color, wrap, draw_lines)

        _callout_pdf(c, d, wrap, draw_lines, W, accent, ink, Color)

        c.setFont("Helvetica", 14)
        c.setFillColor(muted)
        c.drawString(90, 36, th["brand_name"])
        draw_logo()
        if layout != "references":
            draw_citation(d.get("sources"))
        c.showPage()

    c.save()
    return buf.getvalue()
